import os
import uuid
from pathlib import Path
from typing import Optional

import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pdf2image import convert_from_path
import pytesseract

from .config import OUTPUT_DIR, OCR_LANG, ENABLE_OCR


def _add_textbox(slide, left, top, width, height, text, font_size=18):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    return tx


def _ocr_page_image(image: Image.Image) -> str:
    return pytesseract.image_to_string(image, lang=OCR_LANG)


def convert_pdf_to_pptx(pdf_path: str, output_path: Optional[str] = None) -> str:
    pdf_path = str(pdf_path)
    if output_path is None:
        output_path = str(OUTPUT_DIR / f"{uuid.uuid4().hex}.pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    doc = fitz.open(pdf_path)
    page_images = convert_from_path(pdf_path, dpi=180)

    for i, page in enumerate(doc):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Extract native text first
        blocks = page.get_text("blocks")
        has_text = False
        for b in blocks:
            x0, y0, x1, y1, text, block_no, block_type = b[:7]
            text = (text or "").strip()
            if not text:
                continue
            has_text = True
            left = Inches((x0 / page.rect.width) * 13.333)
            top = Inches((y0 / page.rect.height) * 7.5)
            width = Inches(((x1 - x0) / page.rect.width) * 13.333)
            height = Inches(((y1 - y0) / page.rect.height) * 7.5)
            _add_textbox(slide, left, top, width, height, text, font_size=14)

        # OCR fallback when no text is found or OCR is enabled for verification
        if ENABLE_OCR and i < len(page_images):
            ocr_text = _ocr_page_image(page_images[i]).strip()
            if ocr_text and not has_text:
                _add_textbox(slide, Inches(0.5), Inches(0.5), Inches(12.0), Inches(6.5), ocr_text, font_size=12)

        # Add page image as background reference if needed for image-heavy PDFs
        # For editable output, keep this off by default.

    if len(prs.slides) > 0:
        # Remove the default first blank slide if present in some environments
        pass

    prs.save(output_path)
    return output_path
